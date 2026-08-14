"""Smart Chapter and Section Splitter service for textbooks and large documents."""
import re
from dataclasses import dataclass
from pathlib import Path
from pydantic import BaseModel

import pymupdf
from docx import Document
from pptx import Presentation

from app.services.ingestion import TextSection


class ChapterItem(BaseModel):
    index: int
    title: str
    start_page: int
    end_page: int
    page_count: int


def detect_chapters(file_path: Path) -> list[ChapterItem]:
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return _detect_pdf_chapters(file_path)
    elif suffix == ".docx":
        return _detect_docx_chapters(file_path)
    elif suffix == ".pptx":
        return _detect_pptx_chapters(file_path)
    return []


def _detect_pdf_chapters(file_path: Path) -> list[ChapterItem]:
    with pymupdf.open(file_path) as doc:
        total_pages = len(doc)
        if total_pages == 0:
            return []

        # 1. Try PyMuPDF Table of Contents (bookmarks)
        toc = doc.get_toc()
        # Filter top-level items (level 1 or 2)
        top_level = [item for item in toc if item[0] <= 2 and item[2] > 0]

        if len(top_level) >= 2:
            chapters: list[ChapterItem] = []
            for i, item in enumerate(top_level):
                title = item[1].strip()
                start_page = max(1, min(item[2], total_pages))
                if i + 1 < len(top_level):
                    end_page = max(start_page, min(top_level[i + 1][2] - 1, total_pages))
                else:
                    end_page = total_pages

                page_count = (end_page - start_page) + 1
                chapters.append(
                    ChapterItem(
                        index=i + 1,
                        title=title or f"Chapter {i + 1}",
                        start_page=start_page,
                        end_page=end_page,
                        page_count=page_count,
                    )
                )
            return chapters

        # 2. Fallback: Regex scan on first lines of each page for Chapter headings
        detected_points: list[tuple[str, int]] = []
        chapter_regex = re.compile(
            r"(?:^|\n)(?:CHAPTER|Chapter|SECTION|Section|UNIT|Unit|MODULE|Module)\s+([0-9IVXLCDM]+[:.\s–-]+[^\n]{2,80})",
            re.IGNORECASE,
        )

        for page_idx, page in enumerate(doc):
            text = page.get_text("text")[:500]  # Check top of page
            match = chapter_regex.search(text)
            if match:
                clean_title = re.sub(r"\s+", " ", match.group(0)).strip()
                detected_points.append((clean_title, page_idx + 1))

        if len(detected_points) >= 2:
            chapters = []
            for i, (title, start_page) in enumerate(detected_points):
                if i + 1 < len(detected_points):
                    end_page = max(start_page, detected_points[i + 1][1] - 1)
                else:
                    end_page = total_pages
                chapters.append(
                    ChapterItem(
                        index=i + 1,
                        title=title,
                        start_page=start_page,
                        end_page=end_page,
                        page_count=(end_page - start_page) + 1,
                    )
                )
            return chapters

        # 3. Fallback: If document is large (>30 pages), split into logical parts of 25 pages
        if total_pages > 30:
            part_size = 25
            chapters = []
            part_idx = 1
            for start in range(1, total_pages + 1, part_size):
                end = min(start + part_size - 1, total_pages)
                chapters.append(
                    ChapterItem(
                        index=part_idx,
                        title=f"Section {part_idx} (Pages {start}–{end})",
                        start_page=start,
                        end_page=end,
                        page_count=(end - start) + 1,
                    )
                )
                part_idx += 1
            return chapters

        # Single chapter
        return [
            ChapterItem(
                index=1,
                title=Path(file_path).stem,
                start_page=1,
                end_page=total_pages,
                page_count=total_pages,
            )
        ]


def _detect_docx_chapters(file_path: Path) -> list[ChapterItem]:
    doc = Document(file_path)
    headings = []
    for i, p in enumerate(doc.paragraphs):
        if p.style.name.startswith("Heading 1") and p.text.strip():
            headings.append(p.text.strip())

    if len(headings) >= 2:
        return [
            ChapterItem(
                index=idx + 1,
                title=title,
                start_page=idx + 1,
                end_page=idx + 1,
                page_count=1,
            )
            for idx, title in enumerate(headings)
        ]

    return [
        ChapterItem(
            index=1,
            title=Path(file_path).stem,
            start_page=1,
            end_page=1,
            page_count=1,
        )
    ]


def _detect_pptx_chapters(file_path: Path) -> list[ChapterItem]:
    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    if total_slides <= 30:
        return [
            ChapterItem(
                index=1,
                title=Path(file_path).stem,
                start_page=1,
                end_page=total_slides,
                page_count=total_slides,
            )
        ]

    # Split slides into 25-slide modules
    part_size = 25
    chapters = []
    part_idx = 1
    for start in range(1, total_slides + 1, part_size):
        end = min(start + part_size - 1, total_slides)
        chapters.append(
            ChapterItem(
                index=part_idx,
                title=f"Module {part_idx} (Slides {start}–{end})",
                start_page=start,
                end_page=end,
                page_count=(end - start) + 1,
            )
        )
        part_idx += 1
    return chapters


def extract_chapter_sections(
    file_path: Path, start_page: int, end_page: int
) -> list[TextSection]:
    """Extract text sections specifically for the requested page range."""
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        sections: list[TextSection] = []
        with pymupdf.open(file_path) as doc:
            for page_num in range(start_page, min(end_page + 1, len(doc) + 1)):
                page = doc[page_num - 1]
                text = page.get_text("text")
                if text.strip():
                    sections.append(TextSection(text=text, page_number=page_num))
        return sections

    elif suffix == ".docx":
        doc = Document(file_path)
        text = "\n".join(
            p.text for p in doc.paragraphs if p.text.strip()
        )
        return [TextSection(text=text, page_number=1)] if text else []

    elif suffix == ".pptx":
        prs = Presentation(file_path)
        sections = []
        for idx in range(start_page, min(end_page + 1, len(prs.slides) + 1)):
            slide = prs.slides[idx - 1]
            text = "\n".join(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            if text:
                sections.append(TextSection(text=text, page_number=idx))
        return sections

    return []
