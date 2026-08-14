import re
from dataclasses import dataclass
from pathlib import Path

import pymupdf
from docx import Document
from pptx import Presentation


ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".docx"}


@dataclass
class TextSection:
    text: str
    page_number: int


@dataclass
class Chunk:
    content: str
    page_number: int
    chunk_index: int


def extract_text(file_path: Path) -> list[TextSection]:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        with pymupdf.open(file_path) as document:
            return [TextSection(page.get_text("text"), index + 1) for index, page in enumerate(document) if page.get_text("text").strip()]
    if suffix == ".docx":
        document = Document(file_path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        return [TextSection(text, 1)] if text else []
    if suffix == ".pptx":
        presentation = Presentation(file_path)
        sections = []
        for index, slide in enumerate(presentation.slides):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip())
            if text:
                sections.append(TextSection(text, index + 1))
        return sections
    raise ValueError("Only PDF, PPTX, and DOCX files are supported.")


def chunk_sections(sections: list[TextSection], target_words: int = 380, overlap_words: int = 60) -> list[Chunk]:
    chunks: list[Chunk] = []
    chunk_index = 0
    for section in sections:
        normalized = re.sub(r"\s+", " ", section.text).strip()
        words = normalized.split()
        start = 0
        while start < len(words):
            end = min(start + target_words, len(words))
            content = " ".join(words[start:end])
            if content:
                chunks.append(Chunk(content=content, page_number=section.page_number, chunk_index=chunk_index))
                chunk_index += 1
            if end == len(words):
                break
            start = end - overlap_words
    return chunks
